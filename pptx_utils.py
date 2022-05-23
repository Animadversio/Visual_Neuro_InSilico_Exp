import pptx
from pptx import Presentation
from pptx.util import Inches, Length, Pt
from os.path import join
from tqdm import tqdm
import numpy as np

def view_layout_params(pptx_path, slides_num=1, digits=3):
    pprs = Presentation(pptx_path)
    layout_dict = {}
    for shape in pprs.slides[slides_num].shapes:
        print("Object:", shape.name)
        print(type(shape))
        pos_dict = {"height" : np.around(Length(shape.height).inches, digits),
                    "width" : np.around(Length(shape.width).inches, digits),
                    "top" : np.around(Length(shape.top).inches, digits),
                    "left" : np.around(Length(shape.left).inches, digits),}
        print(pos_dict)
        layout_dict[shape.name] = pos_dict
        if hasattr(shape,"text"):
            print("Text ", shape.text)

        if isinstance(shape, pptx.shapes.picture.Picture):
            crop_dict={"crop_right" : np.around(shape.crop_right, digits),
                        "crop_left" : np.around(shape.crop_left, digits),
                        "crop_top" : np.around(shape.crop_top, digits),
                        "crop_bottom" : np.around(shape.crop_bottom, digits), }
            print(crop_dict)
            layout_dict[shape.name].update(crop_dict)
    return layout_dict
