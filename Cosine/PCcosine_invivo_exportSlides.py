import pptx
from pptx import Presentation
from pptx.util import Inches, Length, Pt
from tqdm import tqdm
import os
from os.path import join
from pptx_utils import *
figroot = r"E:\OneDrive - Washington University in St. Louis\Evol_PCCosine"
sumdir = r"E:\OneDrive - Washington University in St. Louis\Evol_PCCosine\summarySlides"
#%%
expdirs = os.listdir(figroot)
expdirs = [fdr for fdr in expdirs if not ("refRepr" in fdr or "RF" in fdr or "summary" in fdr)]
print("%d PC Cosine Experiments from %s to %s"%
      (len(expdirs), expdirs[0][:10], expdirs[-1][:10]))
#%%
def layout_img_on_slide(prs, img_root, explabel):
    protodim1 = Inches(2.05)
    protodim2 = Inches(2.65)

    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    tf = slide.shapes.title
    tf.text = explabel
    tf.top = Inches(0.0)
    tf.left = Inches(0.0)
    tf.width = Inches(8.6875)
    tf.height = protodim1
    tf.text_frame._set_font("Candara", 48, False, False)

    pic = slide.shapes.add_picture(join(img_root, "NormAct_Pattern_TargArea_Evolution.png"),
                                   left=Inches(0), top=protodim1, width=Inches(8.6875), height=Inches(5.1865))
    pic.crop_right = 0.09349
    pic.crop_left = 0.07378
    pic2 = slide.shapes.add_picture(join(img_root, "score_offline_scoretraj.png"),
                                    left=Inches(8.6875), top=Inches(2.5895), width=Inches(4.6458),
                                    height=Inches(4.4845))
    pic2.crop_right = 0.06951
    slide.shapes.add_picture(join(img_root, "alpha_mask_rect_wt_thr.png"),
                             left=Inches(10.96), top=Inches(0), width=protodim1, height=protodim1)

    mtgcrop_setting = {'crop_right': 0.06913, 'crop_left': 0.08617, 'crop_top': 0.0, 'crop_bottom': 0.08611}
    natmtgcrop_setting = {'crop_right': 0.03457, 'crop_left': 0.06283, 'crop_top': 0.0, 'crop_bottom': 0.13194}

    slide2 = prs.slides.add_slide(blank_slide_layout)
    pic3 = slide2.shapes.add_picture(join(img_root, "Image_Evol_per_gen_score_framed.jpg"),
                                     left=Inches(0.0), top=Inches(0), width=Inches(7.7560), height=Inches(7.3325))
    for k, v in mtgcrop_setting.items(): setattr(pic3, k, v)

    slide3 = prs.slides.add_slide(blank_slide_layout)
    pic3 = slide3.shapes.add_picture(join(img_root, "Image_Evol_per_gen_score_framed_wRFmsk_rect_wt_thr.jpg"),
                                     left=Inches(0.0), top=Inches(0), width=Inches(7.7560), height=Inches(7.3325))
    for k, v in mtgcrop_setting.items(): setattr(pic3, k, v)
    slide3.shapes.add_picture(join(img_root, "alpha_mask_rect_wt_thr.png"),
                              left=Inches(9.0), top=Inches(0), width=protodim2, height=protodim2)

    slide4 = prs.slides.add_slide(blank_slide_layout)
    pic3 = slide4.shapes.add_picture(join(img_root, "best_align_natimg_cosine_score_framed.jpg"),
                 height=Inches(5.007), width=Inches(6.65), top=Inches(1.59375), left=Inches(0.0), )
    for k, v in natmtgcrop_setting.items(): setattr(pic3, k, v)
    pic4 = slide4.shapes.add_picture(join(img_root, "best_align_natimg_cosine_neg_score_framed.jpg"),
                 height=Inches(5.007), width=Inches(6.65), top=Inches(1.59375), left=Inches(6.65), )
    for k, v in natmtgcrop_setting.items(): setattr(pic4, k, v)
    txBox = slide4.shapes.add_textbox(Inches(0.0), Inches(1.0), Inches(6.65), Inches(0.4))
    txBox.text_frame.text = "Positive dir alignment"
    txBox = slide4.shapes.add_textbox(Inches(6.65), Inches(1.0), Inches(6.65), Inches(0.4))
    txBox.text_frame.text = "Negative dir alignment"

    slide5 = prs.slides.add_slide(blank_slide_layout)
    pic3 = slide5.shapes.add_picture(join(img_root, "best_align_natimg_cosine_score_framed_wRFmsk_rect_wt_thr.jpg"),
                 height=Inches(5.007), width=Inches(6.65), top=Inches(1.59375), left=Inches(0.0), )
    for k, v in natmtgcrop_setting.items(): setattr(pic3, k, v)
    pic3 = slide5.shapes.add_picture(join(img_root, "best_align_natimg_cosine_neg_score_framed_wRFmsk_rect_wt_thr.jpg"),
                 height=Inches(5.007), width=Inches(6.65), top=Inches(1.59375), left=Inches(6.65), )
    for k, v in natmtgcrop_setting.items(): setattr(pic3, k, v)
    txBox = slide5.shapes.add_textbox(Inches(0.0), Inches(1.0), Inches(6.65), Inches(0.4))
    txBox.text_frame.text = "Positive dir alignment"
    txBox = slide5.shapes.add_textbox(Inches(6.65), Inches(1.0), Inches(6.65), Inches(0.4))
    txBox.text_frame.text = "Negative dir alignment"
    # slide5.shapes.add_picture(join(img_root, "alpha_mask_rect_wt_thr.png"),
    #                           left=Inches(9.0), top=Inches(0), width=protodim2, height=protodim2)
    return

prs = Presentation()
# 16:9 wide screen layout
prs.slide_width = Inches(13.33333)
prs.slide_height = Inches(7.5)
for expdir in expdirs:
    img_root = join(figroot, expdir)
    layout_img_on_slide(prs, img_root, expdir)


prs.save(join(sumdir, "PCCosine_Alfa_summary.pptx"))
#%%
"alpha_mask_rect_wt_thr.png",
"NormAct_Pattern_TargArea_Evolution.png",
'score_offline_scoretraj.png',
"Image_Evol_per_gen_score_framed_wRFmsk_rect_wt_thr.jpg",
"Image_Evol_per_gen_score_framed.jpg",
"best_align_natimg_cosine_score_framed_wRFmsk_rect_wt_thr.jpg",
"best_align_natimg_cosine_score_framed.jpg"
#%%

view_layout_params(join(sumdir, "demo.pptx"), slides_num=0)
