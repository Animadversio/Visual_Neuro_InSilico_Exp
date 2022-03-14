#%%
#%%
import pptx
from pptx import Presentation
from pptx.util import Inches, Length, Pt
from tqdm import tqdm
from os.path import join
from pptx_utils import *

figroot = r"E:\OneDrive - Washington University in St. Louis\Evol_Cosine"
sumdir = r"E:\OneDrive - Washington University in St. Louis\Evol_Cosine\summarySlides"

#%%
view_layout_params(join(sumdir, "Demo.pptx"), slides_num=0)
view_layout_params(join(sumdir, "Demo.pptx"), slides_num=1)
view_layout_params(join(sumdir, "Demo.pptx"), slides_num=2)
#%% slides 1
def layout_img_on_slide(prs, imgroot, explabel):
    protodim1 = Inches(2.05)
    protodim2 = Inches(2.65)
    title_setting = {"top": Inches(0.0), "left": Inches(0.0),
                     "width": Inches(8.6875), "height": protodim1,}
    blank_slide_layout = prs.slide_layouts[5]
    # First slide: Population evolution patterns; Score trajectory; Target images; Masks
    slide = prs.slides.add_slide(blank_slide_layout)
    tf = slide.shapes.title
    tf.text = expdir
    tf.text_frame._set_font("Candara", 48, False, False)
    for k, v in title_setting.items(): setattr(tf, k, v)
    pic = slide.shapes.add_picture(join(img_root, "NormAct_Pattern_TargArea_Evolution.png"),
                                   left=Inches(0), top=protodim1, width=Inches(8.6875), height=Inches(5.1865))
    pic.crop_right = 0.09349
    pic.crop_left = 0.07378
    pic2 = slide.shapes.add_picture(join(img_root, "score_offline_scoretraj.png"),
                                    left=Inches(8.6875), top=Inches(2.5895), width=Inches(4.6458),
                                    height=Inches(4.4845))
    pic2.crop_right = 0.06951
    slide.shapes.add_picture(join(img_root, "targetimg.png"),
                             left=Inches(8.88), top=Inches(0), width=protodim1, height=protodim1)
    slide.shapes.add_picture(join(img_root, "targetimg.png"),
                             left=Inches(10.96), top=Inches(0), width=protodim1, height=protodim1)
    slide.shapes.add_picture(join(img_root, "alpha_mask_rect_wt_thr.png"),
                             left=Inches(10.96), top=Inches(0), width=protodim1, height=protodim1)
    # Second slide: Un-normalized Population evolution patterns; Raw Score points; Target images; Masks
    slide = prs.slides.add_slide(blank_slide_layout)
    tf = slide.shapes.title
    tf.text = expdir
    tf.text_frame._set_font("Candara", 48, False, False)
    for k, v in title_setting.items(): setattr(tf, k, v)
    pic = slide.shapes.add_picture(join(img_root, "RawAct_Pattern_TargArea_Evolution.png"),
                                   left=Inches(0), top=protodim1, width=Inches(8.6875), height=Inches(5.1865))
    pic.crop_right = 0.09349
    pic.crop_left = 0.07378
    pic2 = slide.shapes.add_picture(join(img_root, "score_offline_scoretraj_w_pnts.png"),
                                    left=Inches(8.6875), top=Inches(2.5895), width=Inches(4.6458),
                                    height=Inches(4.4845))
    pic2.crop_right = 0.06951
    slide.shapes.add_picture(join(img_root, "targetimg.png"),
                             left=Inches(8.88), top=Inches(0), width=protodim1, height=protodim1)
    slide.shapes.add_picture(join(img_root, "targetimg.png"),
                             left=Inches(10.96), top=Inches(0), width=protodim1, height=protodim1)
    slide.shapes.add_picture(join(img_root, "alpha_mask_rect_wt_thr.png"),
                             left=Inches(10.96), top=Inches(0), width=protodim1, height=protodim1)
    # 3rd Slide, Image evolution ; target image
    mtgcrop_setting = {'crop_right': 0.06913, 'crop_left': 0.08617, 'crop_top': 0.0, 'crop_bottom': 0.08611}
    slide2 = prs.slides.add_slide(blank_slide_layout)
    pic3 = slide2.shapes.add_picture(join(img_root, "Image_Evol_per_gen_score_framed.jpg"),
                                     left=Inches(0.0), top=Inches(0), width=Inches(7.7560), height=Inches(7.3325))
    for k, v in mtgcrop_setting.items(): setattr(pic3, k, v)
    slide2.shapes.add_picture(join(img_root, "targetimg.png"),
                              left=Inches(9.0), top=Inches(0), width=protodim2, height=protodim2)
    # 4th Slide, Image evolution with RF mask ; target image with RF mask
    slide3 = prs.slides.add_slide(blank_slide_layout)
    pic3 = slide3.shapes.add_picture(join(img_root, "Image_Evol_per_gen_score_framed_wRFmsk_rect_wt_thr.jpg"),
                                     left=Inches(0.0), top=Inches(0), width=Inches(7.7560), height=Inches(7.3325))
    for k, v in mtgcrop_setting.items(): setattr(pic3, k, v)
    slide3.shapes.add_picture(join(img_root, "targetimg.png"),
                              left=Inches(9.0), top=Inches(0), width=protodim2, height=protodim2)
    slide3.shapes.add_picture(join(img_root, "alpha_mask_rect_wt_thr.png"),
                              left=Inches(9.0), top=Inches(0), width=protodim2, height=protodim2)
    return

expdirs = ["2022-03-02-Beto-01-cosine_V1V4",
          "2022-03-02-Beto-02-cosine_V1",
          "2022-03-02-Beto-03-cosine_V1",
          "2022-03-02-Beto-04-cosine_V1V4",]
#%%
import os
expdirs = os.listdir(figroot)
expdirs = [fdr for fdr in expdirs if not ("refRepr" in fdr or "RF" in fdr or "summary" in fdr) and  ("Beto" in fdr)]
print("%d Cosine Experiments from %s to %s"%
      (len(expdirs), expdirs[0][:10], expdirs[-1][:10]))
#%%
prs = Presentation()
# 16:9 wide screen layout
prs.slide_width = Inches(13.33333)
prs.slide_height = Inches(7.5)
for expdir in expdirs:
    img_root = join(figroot, expdir)
    layout_img_on_slide(prs, img_root, expdir)

prs.save(join(sumdir, "CosineEvol_Beto_summary.pptx"))
#%%

#%%
"score_offline_scoretraj.png"
"target.png"
"RawAct_Pattern_TargArea_Evolution.png"
"Image_Evol_per_gen_score_framed_wRFmsk_rect_wt_thr.png"