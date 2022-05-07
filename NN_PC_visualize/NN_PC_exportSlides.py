""" 
Merge the figures into a single slides show
"""
from pptx import Presentation
from pptx.util import Inches, Length, Pt
from os.path import join
from tqdm import tqdm
from NN_PC_visualize.NN_PC_lib import shorten_layername
outdir = r"H:\CNN-PCs\summary"

def layout_img_on_slide(slide, titlestr, negvecnm, posvecnm, with_rfmsk=False, img_root=None):
    tf = slide.shapes.title
    tf.text = titlestr
    tf.top = Inches(0.0)
    tf.width = Inches(13.3)
    tf.height = Inches(1.8)
    tf.text_frame._set_font("Candara", 48, False, False)

    dir = "neg"
    # vecnm = "%s_%s_PC%03d_%s_cosine" % (netname, shorten_layername(layer), iPC, dir)
    width = Inches(3.2)
    left1 = Inches(0.0); top1 = Inches(2)
    left2 = Inches(3.4); top2 = Inches(3.8)
    img_path = join(img_root, f"{negvecnm}_evolpatch%s.png"%('_w_rfmsk' if with_rfmsk else ''))
    pic = slide.shapes.add_picture(img_path, left2, top1, width=width)
    img_path = join(img_root, f"{negvecnm}_natrpatch%s.png"%('_w_rfmsk' if with_rfmsk else ''))
    pic = slide.shapes.add_picture(img_path, left2, top2, width=width)
    img_path = join(img_root, f"{negvecnm}_scoretraj.png")
    pic = slide.shapes.add_picture(img_path, left1, top1, width=width)
    pic.crop_right = 0.098
    pic.height = Length(pic.height / (1 - 0.098))

    txBox = slide.shapes.add_textbox(left2, top2 + width, width, Inches(1.0))
    txBox.text_frame.text = "Negative dir"

    dir = "pos"
    # vecnm = "%s_%s_PC%03d_%s_cosine" % (netname, shorten_layername(layer), iPC, dir)
    width = Inches(3.2)
    left1 = Inches(10.13333); top1 = Inches(2)
    left2 = Inches(6.8); top2 = Inches(3.8)
    img_path = join(img_root, f"{posvecnm}_evolpatch%s.png"%('_w_rfmsk' if with_rfmsk else ''))
    pic = slide.shapes.add_picture(img_path, left2, top1, width=width)
    img_path = join(img_root, f"{posvecnm}_natrpatch%s.png"%('_w_rfmsk' if with_rfmsk else ''))
    pic = slide.shapes.add_picture(img_path, left2, top2, width=width)
    img_path = join(img_root, f"{posvecnm}_scoretraj.png")
    pic = slide.shapes.add_picture(img_path, left1, top1, width=width)
    pic.crop_right = 0.098
    pic.height = Length(pic.height / (1 - 0.098))
    txBox = slide.shapes.add_textbox(left2, top2 + width, width, Inches(1.0))
    txBox.text_frame.text = "Positive dir"


def export_imgs_to_ppt(netname, reclayers, img_root, save_prefix="", rfmask_slide=False, fullimg_slide=True):
    for layeri, layer in enumerate(reclayers):
        prs = Presentation()
        # 16:9 wide screen layout
        prs.slide_width = Inches(13.33333)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[5]
        print("Processing layer %s (%d/%d)"%(layer, layeri, len(reclayers)))
        for iPC in tqdm(range(100)):
            titlestr = "%s %s PC%03d (Cosine)" % (netname, shorten_layername(layer), iPC)

            negvecnm = "%s_%s_PC%03d_%s_cosine" % (netname, shorten_layername(layer), iPC, "neg")
            posvecnm = "%s_%s_PC%03d_%s_cosine" % (netname, shorten_layername(layer), iPC, "pos")
            if fullimg_slide:
                slide = prs.slides.add_slide(blank_slide_layout)
                layout_img_on_slide(slide, titlestr, negvecnm, posvecnm,
                                    with_rfmsk=False, img_root=img_root)
            if rfmask_slide:
                slide = prs.slides.add_slide(blank_slide_layout)
                layout_img_on_slide(slide, titlestr, negvecnm, posvecnm,
                                    with_rfmsk=True, img_root=img_root)

        prs.save(join(outdir, '%s_%s_PCdir_improve_visualize.pptx'%(save_prefix, shorten_layername(layer))))

#%%

#%%
netname = "resnet50_linf8"
save_prefix = "ResNetRobust"
reclayers = ['.layer2.Bottleneck2', '.layer3.Bottleneck2', '.layer4.Bottleneck2']
# img_root = r"H:\CNN-PCs\natpatch_norm_vis"
# img_root = r"H:\CNN-PCs\natpatch_norm_lr_vis"
img_root = r"H:\CNN-PCs\resnet_linf8_natpatch_norm_lr_reldir_vis"
export_imgs_to_ppt(netname, reclayers, img_root, save_prefix, rfmask_slide=True, fullimg_slide=True)

#%%
netname = "vgg16"
save_prefix = "VGG16"
reclayers = ['.features.ReLU15','.features.ReLU29','.classifier.ReLU1','.classifier.ReLU4']
img_root = r"H:\CNN-PCs\VGG16_natpatch_norm_lr_reldir_vis"
export_imgs_to_ppt(netname, reclayers, img_root, save_prefix, rfmask_slide=True, fullimg_slide=True)

#%%
netname = "densenet_robust"
save_prefix = "DenseNet_robust"
reclayers = ['.features._DenseBlockdenseblock1',
             '.features._DenseBlockdenseblock2',
             '.features._DenseBlockdenseblock3',
             '.features._DenseBlockdenseblock4']
img_root = r"H:\CNN-PCs\densenet_robust_natpatch_norm_lr_reldir_vis"
export_imgs_to_ppt(netname, reclayers, img_root, save_prefix, rfmask_slide=True, fullimg_slide=True)

#%% Note there will be error if you try to open a slide pptx file when it's opened in PPT
from pptx import Presentation
prs2 = Presentation(pptx=join(outdir, 'test.pptx'))
#%%
imgshapes = list(prs2.slides[0].shapes)
for shp in imgshapes:
    print("%.3f %.3f @ (%.3f %.3f) cropright %.4f" %
          (shp.height.inches, shp.width.inches, shp.left.inches, shp.top.inches, shp.crop_right))
