import pptx
from pptx import Presentation
from pptx.util import Inches, Length, Pt
from os.path import join
from tqdm import tqdm

protodir = r"F:\insilico_exps\CorNet-recurrent-evol\proto_summary"
tracedir = r"F:\insilico_exps\CorNet-recurrent-evol\actdyn_summary"

def _layout_figures_on_slide(slide, tracepath, hmappath, protopath, chanlabel):
    tf = slide.shapes.title
    tf.text = chanlabel
    tf.top = Inches(0.0)
    tf.left = Inches(0.0)
    tf.width = Inches(6.1)
    tf.height = Inches(1.6)
    tf.text_frame._set_font("Candara", 48, False, False)
    pic = slide.shapes.add_picture(protopath, Inches(0.25), Inches(1.75), width=Inches(5.6))
    pic = slide.shapes.add_picture(tracepath, Inches(6.45), Inches(0.55), width=Inches(6.16))
    pic = slide.shapes.add_picture(hmappath, Inches(6.10), Inches(5.20), width=Inches(7.24), height=Inches(2.3))
    pic.crop_right = 0.31947
    pic.crop_left = 0.16593
    pic.crop_top = 0.17234
    pic.crop_bottom = 0.04221

def _layout_figures_on_slide2(slide, tracepath, hmappath, protopath, chanlabel):
    tf = slide.shapes.title
    tf.text = chanlabel
    tf.top = Inches(0.0)
    tf.left = Inches(0.0)
    tf.width = Inches(6.1)
    tf.height = Inches(1.6)
    tf.text_frame._set_font("Candara", 48, False, False)
    pic = slide.shapes.add_picture(protopath, Inches(0.25), Inches(1.75), width=Inches(5.6))
    pic = slide.shapes.add_picture(tracepath, Inches(6.45), Inches(0.55), width=Inches(6.16))
    pic = slide.shapes.add_picture(hmappath, Inches(6.85), Inches(5.17), width=Inches(5.6081), height=Inches(2.33))
    pic.crop_right = 0.28132
    pic.crop_left = 0.16593
    pic.crop_top = 0.31779
    pic.crop_bottom = 0.13105

def merge_plots2slides(area, sublayer, chanrng, outdir, save_sfx=""):
    prs = Presentation()
    # 16:9 wide screen layout
    prs.slide_width = Inches(13.33333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[5]
    for channum in tqdm(range(chanrng[0], chanrng[1])):
        chanlabel = f"{area}-{sublayer}-Ch{channum:03d}"
        tracepath = (join(tracedir, f"{chanlabel}_act_traces.png"))
        hmappath = (join(tracedir, f"{chanlabel}_act_heatmap.png"))
        protopath = (join(protodir, f"{chanlabel}_allproto_mtg.jpg"))
        slide = prs.slides.add_slide(blank_slide_layout)
        if area is "V4":
            _layout_figures_on_slide(slide, tracepath, hmappath, protopath, chanlabel)
        elif area in ["V2", "IT"]:
            _layout_figures_on_slide2(slide, tracepath, hmappath, protopath, chanlabel)

    prs.save(join(outdir, f'CorNet-s_{area}-{sublayer}-Ch{chanrng[0]:d}-{chanrng[1]:d}_evol_dynam_visualize{save_sfx}.pptx'))
#%%
sumdir = r"F:\insilico_exps\CorNet-recurrent-evol\summary"
merge_plots2slides("V4", "output", (0, 100), sumdir, save_sfx="")
#%%
merge_plots2slides("V2", "output", (0, 50), sumdir, save_sfx="")
merge_plots2slides("IT", "output", (0, 100), sumdir, save_sfx="")


#%%
pprs = Presentation(join(sumdir,"CorNet-s_V2-output-Ch0-50_evol_dynam_visualize.pptx"))
for shape in pprs.slides[0].shapes:
    pos_dict = {"height" : Length(shape.height).inches,
                "width" : Length(shape.width).inches,
                "top" : Length(shape.top).inches,
                "left" : Length(shape.left).inches,}
    print(pos_dict)
    if isinstance(shape, pptx.shapes.picture.Picture):
        crop_dict={"crop_right" : shape.crop_right,
        "crop_left" : shape.crop_left,
        "crop_top" : shape.crop_top,
        "crop_bottom" : shape.crop_bottom,}
        print(crop_dict)
